#!/usr/bin/env python3
"""Build AgentMesh.pptx — 11-slide deck, dark 'settlement terminal' theme.
Imports cleanly into Google Slides (Drive > open with Slides, or Slides >
File > Import slides). Fonts chosen from the Google Slides catalogue so they
survive the import: Roboto Mono (display/labels) + Roboto (body)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (dark theme, matches the HTML deck) ----
INK      = RGBColor(0x0A, 0x0F, 0x16)
PANEL    = RGBColor(0x11, 0x19, 0x26)
PANEL2   = RGBColor(0x0E, 0x15, 0x20)
LINE     = RGBColor(0x24, 0x32, 0x44)
HI       = RGBColor(0xE9, 0xEE, 0xF4)
MID      = RGBColor(0x9B, 0xA9, 0xBC)
LOW      = RGBColor(0x6B, 0x7A, 0x8E)
JADE     = RGBColor(0x2F, 0xD0, 0xA0)
BLUE     = RGBColor(0x4C, 0x8D, 0xFF)
GOLD     = RGBColor(0xE5, 0xB5, 0x67)

MONO = "Roboto Mono"
SANS = "Roboto"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

MX = Inches(0.85)          # left margin
CW = EMU_W - MX * 2        # content width


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def _set_run(r, text, size, color, font=SANS, bold=False, spacing=None, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if spacing is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))


def textbox(s, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def eyebrow(s, num, label):
    tb, tf = textbox(s, MX, Inches(0.62), CW, Inches(0.4))
    p = tf.paragraphs[0]
    _set_run(p.add_run(), num + "   ", 12, JADE, MONO, bold=True, spacing=2.2)
    _set_run(p.add_run(), label.upper(), 12, LOW, MONO, spacing=2.2)
    # hairline
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, Inches(1.02), CW, Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background()
    ln.shadow.inherit = False
    return tb


def heading(s, top, parts, size=40):
    """parts: list of (text, color) tuples; supports \\n."""
    tb, tf = textbox(s, MX, top, CW, Inches(1.7))
    p = tf.paragraphs[0]
    p.line_spacing = 1.02
    for text, color in parts:
        segs = text.split("\n")
        for j, seg in enumerate(segs):
            if j > 0:
                p = tf.add_paragraph(); p.line_spacing = 1.02
            if seg:
                _set_run(p.add_run(), seg, size, color, MONO, bold=True, spacing=-0.4)
    return tb


def lead(s, top, runs, width=None, size=17):
    tb, tf = textbox(s, MX, top, width or Inches(8.6), Inches(1.4))
    p = tf.paragraphs[0]; p.line_spacing = 1.4
    for text, color, bold in runs:
        _set_run(p.add_run(), text, size, color, SANS, bold=bold)
    return tb


def rrect(s, left, top, width, height, fill=PANEL, line=LINE):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.045
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def left_align(tf):
    """Auto-shape text frames inherit centered paragraphs; force left."""
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT


# =====================================================================
# 1 · TITLE
# =====================================================================
s = slide()
tb, tf = textbox(s, MX, Inches(0.62), CW, Inches(0.4))
p = tf.paragraphs[0]
_set_run(p.add_run(), "◆  ", 12, JADE, MONO, spacing=2.2)
_set_run(p.add_run(), "ENCODE · PROGRAMMABLE MONEY HACKATHON — AGENTIC ECONOMY", 12, LOW, MONO, spacing=2.2)

tb, tf = textbox(s, MX, Inches(2.05), CW, Inches(2.0))
p = tf.paragraphs[0]
_set_run(p.add_run(), "Agent", 96, HI, MONO, bold=True, spacing=-1)
_set_run(p.add_run(), "Mesh", 96, JADE, MONO, bold=True, spacing=-1)

lead(s, Inches(4.15), [
    ("The compliant settlement layer for the agent economy. ", MID, False),
    ("Named wallets, sub-cent payments, and screened escrow", HI, True),
    (" — every leg settled in USDC on Arc.", MID, False),
], width=Inches(10.5), size=19)

tb, tf = textbox(s, MX, Inches(5.55), CW, Inches(0.4))
p = tf.paragraphs[0]
tick = "USDC-native      Arc L1 · chain 5042002      x402 micropayments      Circle Agent Stack"
_set_run(p.add_run(), tick, 13, LOW, MONO, spacing=0.6)

pill = rrect(s, MX, Inches(6.35), Inches(5.9), Inches(0.55), fill=PANEL2, line=JADE)
pill.text_frame.word_wrap = True
pp = pill.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
_set_run(pp.add_run(), "●  Working end-to-end on Arc rails today", 13, JADE, MONO, spacing=0.5)


# =====================================================================
# 2 · PROBLEM
# =====================================================================
s = slide()
eyebrow(s, "01", "The problem")
heading(s, Inches(1.35), [("Agents can already spend.\n", HI), ("Nothing lets them ", HI), ("settle", JADE), (".", HI)], size=40)

cards = [
    ("NO IDENTITY", "An address, not an account",
     "Agents have a hex address — no name, no discoverable service, no reputation, no way to prove who they are before money moves."),
    ("MICRO DOESN'T FIT", "Card rails floor at ~$0.30",
     "An agent buying a $0.002 datapoint per call can't touch card networks — the minimum fee is 150× the purchase."),
    ("TRUST GAP", "Nobody wires $5k blind",
     "No one lets an autonomous agent send five figures to an unscreened counterparty with no escrow and no recourse."),
]
cw = (CW - Inches(0.4)) / 3
for i, (k, t, d) in enumerate(cards):
    left = MX + i * (cw + Inches(0.2))
    box = rrect(s, left, Inches(3.35), cw, Inches(3.2))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22); tf.margin_top = Inches(0.22)
    p = tf.paragraphs[0]; _set_run(p.add_run(), k, 12, JADE, MONO, bold=True, spacing=1.4)
    p = tf.add_paragraph(); p.space_before = Pt(10); _set_run(p.add_run(), t, 18, HI, SANS, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(8); p.line_spacing = 1.4
    _set_run(p.add_run(), d, 13.5, MID, SANS)
    left_align(tf)


# =====================================================================
# 3 · SOLUTION (table)
# =====================================================================
s = slide()
eyebrow(s, "02", "The solution")
heading(s, Inches(1.35), [("A ", HI), ("trust stack", JADE), (" for agent money", HI)], size=40)
lead(s, Inches(2.35), [
    ("Stripe + escrow + KYC for the agent economy — six layers, one settlement asset. Arc supplies the rails; AgentMesh supplies the trust.", MID, False)
], width=Inches(11), size=15)

rows = [
    ("Identity", "Named agent wallets — databot.agent.arc + card", "ERC-721 registry · Circle Wallets"),
    ("Payments", "Sub-cent, per-call streaming", "x402 handshake on Arc"),
    ("Settlement", "High-value jobs locked in escrow", "AgentEscrow.sol · USDC"),
    ("Compliance", "Seller screened before funds release", "ComplianceGate · Circle CE"),
    ("Automation", "Auto-release, deadline refunds", "event-driven watcher"),
    ("Brain", "Any MCP model plugs in as an agent", "@agentmesh/mcp-server"),
]
top = Inches(3.25)
tbl_h = Inches(3.7)
gtbl = s.shapes.add_table(len(rows) + 1, 3, MX, top, CW, tbl_h).table
gtbl.columns[0].width = Inches(2.3)
gtbl.columns[1].width = Inches(5.5)
gtbl.columns[2].width = CW - Inches(2.3) - Inches(5.5)
hdr = ["LAYER", "WHAT IT DOES", "POWERED BY"]
for c, h in enumerate(hdr):
    cell = gtbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = PANEL2
    cell.margin_left = Inches(0.16); cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
    p = cell.text_frame.paragraphs[0]; _set_run(p.add_run(), h, 11, LOW, MONO, spacing=1.4)
for r, (a, b, c) in enumerate(rows, start=1):
    for ci, (val, col, fnt, bold) in enumerate([
        (a, HI, MONO, True), (b, MID, SANS, False), (c, BLUE, MONO, False)
    ]):
        cell = gtbl.cell(r, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL if r % 2 else PANEL2
        cell.margin_left = Inches(0.16); cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        _set_run(p.add_run(), val, 12.5 if ci else 13, col, fnt, bold=bold)


# =====================================================================
# 4 · HOW IT WORKS
# =====================================================================
s = slide()
eyebrow(s, "03", "Programmable money flows")
heading(s, Inches(1.35), [("Two settlement modes,\n", HI), ("one ", HI), ("conditional", JADE), (" spine", HI)], size=36)

def flow_panel(left, tag, title, steps):
    box = rrect(s, left, Inches(3.05), (CW - Inches(0.3)) / 2, Inches(3.3))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25); tf.margin_top = Inches(0.22)
    p = tf.paragraphs[0]
    _set_run(p.add_run(), "[" + tag + "]  ", 12, JADE, MONO, bold=True, spacing=1.2)
    _set_run(p.add_run(), title.upper(), 12, LOW, MONO, spacing=1.2)
    for i, st in enumerate(steps):
        p = tf.add_paragraph(); p.space_before = Pt(9); p.line_spacing = 1.3
        _set_run(p.add_run(), str(i + 1) + "  ", 12, JADE, MONO, bold=True)
        _set_run(p.add_run(), st, 13, MID, SANS)
    left_align(tf)

flow_panel(MX, "x402", "Streaming micropayments", [
    "Agent hits a priced endpoint → HTTP 402 + a signed quote",
    "Pays on-chain USDC, signs quoteId‖txHash",
    "Retries with X-PAYMENT; server verifies the transfer on-chain",
    "Payer-bound & replay-safe — observed transfers can't be stolen",
])
flow_panel(MX + (CW - Inches(0.3)) / 2 + Inches(0.3), "escrow", "Multi-step settlement", [
    "createJob — buyer locks USDC for a named seller",
    "deliver — seller submits proof hash on-chain",
    "Dispute window → release, gated on isAllowed(seller)",
    "Branches: dispute→arbiter · deadline→refund · blocked→refundBlocked",
])
lead(s, Inches(6.55), [
    ("Money release is ", MID, False),
    ("conditional on compliance, delivery, and time", HI, True),
    (" — a real multi-step flow, not a bare transfer.", MID, False),
], width=Inches(11.5), size=14)


# =====================================================================
# 5 · ARCHITECTURE
# =====================================================================
s = slide()
eyebrow(s, "04", "Architecture")
heading(s, Inches(1.35), [("One monorepo, ", HI), ("seven", BLUE), (" moving parts", HI)], size=38)

arch = [
    ("contracts/", "Foundry — AgentRegistry, AgentEscrow, ComplianceGate"),
    ("packages/sdk", "viem client, wallet providers (EOA ⇄ Circle), paidFetch"),
    ("packages/shared", "Arc chain config, ABIs, x402 wire types"),
    ("apps/mcp-server", "agent brain — MCP tools any model can drive"),
    ("apps/seller-agent", "DataAgent — x402-priced API + escrow worker"),
    ("apps/watcher", "screening + auto-release / refund automation"),
    ("apps/dashboard", "Next.js live view — agents, payments, escrow"),
    ("apps/extension", "MV3 popup — human oversight, approve / dispute"),
]
colw = (CW - Inches(0.3)) / 2
rh = Inches(0.78)
for i, (path, desc) in enumerate(arch):
    col = i % 2; row = i // 2
    left = MX + col * (colw + Inches(0.3))
    top = Inches(2.75) + row * (rh + Inches(0.18))
    box = rrect(s, left, top, colw, rh)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_run(p.add_run(), path, 13, JADE, MONO, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(3)
    _set_run(p.add_run(), desc, 12, MID, SANS)
    left_align(tf)


# =====================================================================
# 6 · WHY STABLECOIN-NATIVE
# =====================================================================
s = slide()
eyebrow(s, "05", "Why it matters")
heading(s, Inches(1.35), [("What stablecoin-native\n", HI), ("makes possible", JADE)], size=38)

stats = [("50", JADE, "paid API calls settled on-chain in one demo run"),
         ("$0.001", GOLD, "per-call price — 300× below the card floor"),
         ("6 dp", BLUE, "USDC ERC-20 precision for true sub-cent value"),
         ("<1s", JADE, "Arc finality — settlement at machine speed")]
sw = (CW - Inches(0.6)) / 4
for i, (v, col, l) in enumerate(stats):
    left = MX + i * (sw + Inches(0.2))
    box = rrect(s, left, Inches(3.15), sw, Inches(1.75))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]; _set_run(p.add_run(), v, 34, col, MONO, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(8); p.line_spacing = 1.3
    _set_run(p.add_run(), l, 12, MID, SANS)
    left_align(tf)

bullets = [
    ("Machine-speed markets", " — agents pay per request, not per invoice; value flows continuously."),
    ("Compliance as code", " — screening is an on-chain precondition to payout, not paperwork after the fact."),
    ("Agent-to-agent commerce", " — a named seller advertises services and gets paid without a human in the loop."),
]
tb, tf = textbox(s, MX, Inches(5.3), CW, Inches(1.9))
for i, (b, rest) in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(8); p.line_spacing = 1.35
    _set_run(p.add_run(), "▸  ", 15, JADE, MONO, bold=True)
    _set_run(p.add_run(), b, 15, HI, SANS, bold=True)
    _set_run(p.add_run(), rest, 15, MID, SANS)


# =====================================================================
# 7 · CIRCLE MAP
# =====================================================================
s = slide()
eyebrow(s, "06", "Built on Circle")
heading(s, Inches(1.35), [("Deep in the ", HI), ("Agent Stack", JADE)], size=40)

crows = [
    ("USDC on Arc", "live", JADE, "every leg, x402 and escrow, one settlement asset"),
    ("Circle Wallets", "integrated", JADE, "agent-wallet signer via WALLET_PROVIDER=circle, activating with API key"),
    ("Compliance Engine", "integrated", JADE, "watcher screens sellers; labeled fallback when no key"),
    ("Nanopayments", "upgrade path", GOLD, "x402 leg maps directly onto batched nanopayments"),
    ("Paymaster · CCTP", "roadmap", GOLD, "gasless agent onboarding + cross-chain USDC withdrawal"),
]
top = Inches(2.6)
for i, (prod, status, scol, desc) in enumerate(crows):
    box = rrect(s, MX, top + i * Inches(0.86), CW, Inches(0.72))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_run(p.add_run(), prod + "    ", 15, HI, MONO, bold=True)
    mark = "● " if scol == JADE else "◆ "
    _set_run(p.add_run(), mark + status + "   ", 13, scol, MONO, bold=True)
    _set_run(p.add_run(), "— " + desc, 13, MID, SANS)
    left_align(tf)


# =====================================================================
# 8 · TRACK FIT
# =====================================================================
s = slide()
eyebrow(s, "07", "Track fit")
heading(s, Inches(1.35), [("Made for ", HI), ("Agentic Economy", JADE)], size=40)

trows = [
    ("Agents that hold wallets", "named agent wallets — registry + Circle Wallets signer"),
    ("Decision logic on real signals", "watcher acts on on-chain JobCreated / delivery / deadlines"),
    ("Autonomous USDC settlement", "x402 streaming + escrow, no human in the loop"),
    ("Agent Stack → wallets & actions", "SDK wallet layer, MCP tools, on-chain calls"),
    ("Nanopayments for a2a payments", "x402 micropayment handshake (nanopayments-ready)"),
]
gtbl = s.shapes.add_table(len(trows) + 1, 2, MX, Inches(2.65), CW, Inches(3.8)).table
gtbl.columns[0].width = Inches(5.3)
gtbl.columns[1].width = CW - Inches(5.3)
for c, h in enumerate(["WHAT THE TRACK WANTS", "WHAT AGENTMESH SHIPS"]):
    cell = gtbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = PANEL2
    cell.margin_left = Inches(0.16); cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
    _set_run(cell.text_frame.paragraphs[0].add_run(), h, 11, LOW, MONO, spacing=1.2)
for r, (a, b) in enumerate(trows, start=1):
    for ci, (val, col, fnt, bold) in enumerate([(a, HI, SANS, True), (b, BLUE, MONO, False)]):
        cell = gtbl.cell(r, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL if r % 2 else PANEL2
        cell.margin_left = Inches(0.16); cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_run(cell.text_frame.paragraphs[0].add_run(), val, 12.5, col, fnt, bold=bold)


# =====================================================================
# 9 · PRODUCTION-GRADE
# =====================================================================
s = slide()
eyebrow(s, "08", "Not a toy")
heading(s, Inches(1.35), [("Hardened to ", HI), ("production grade", JADE)], size=38)

pstats = [("85", JADE, "tests — 51 Foundry (unit·fuzz·invariant) + 34 vitest"),
          ("97%", BLUE, "contract line coverage"),
          ("CI", JADE, "green — build, typecheck, lint, tests, Slither"),
          ("kill-9", GOLD, "resilience drill passed")]
sw = (CW - Inches(0.6)) / 4
for i, (v, col, l) in enumerate(pstats):
    left = MX + i * (sw + Inches(0.2))
    box = rrect(s, left, Inches(2.55), sw, Inches(1.6))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.16); tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]; _set_run(p.add_run(), v, 30, col, MONO, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(6); p.line_spacing = 1.25
    _set_run(p.add_run(), l, 11.5, MID, SANS)
    left_align(tf)

cols = [
    ("SECURITY", [
        ("Payer-bound x402", " — signed quotes, replay set survives restarts"),
        ("Fail-closed compliance", " — screening outage denies, never opens"),
        ("Ownable2Step + Pausable", " — pause blocks new jobs, never traps funds"),
    ]),
    ("OPERATIONS", [
        ("SQLite persistence", " — payments & job state survive crashes"),
        ("Health checks + logs", " — Docker, graceful shutdown"),
        ("Runbook + deploy guide", " — key rotation, pause, gate-swap"),
    ]),
]
for ci, (title, items) in enumerate(cols):
    left = MX + ci * (CW / 2)
    tb, tf = textbox(s, left, Inches(4.5), CW / 2 - Inches(0.3), Inches(2.6))
    p = tf.paragraphs[0]; _set_run(p.add_run(), title, 12, LOW, MONO, spacing=1.4)
    for b, rest in items:
        p = tf.add_paragraph(); p.space_before = Pt(9); p.line_spacing = 1.3
        _set_run(p.add_run(), "▸  ", 14, JADE, MONO, bold=True)
        _set_run(p.add_run(), b, 14, HI, SANS, bold=True)
        _set_run(p.add_run(), rest, 14, MID, SANS)


# =====================================================================
# 10 · STATUS
# =====================================================================
s = slide()
eyebrow(s, "09", "Where it stands")
heading(s, Inches(1.35), [("Runs today. ", HI), ("Testnet next.", JADE)], size=40)

now = [
    ("pnpm demo:local — full e2e on a local Arc chain in ~2 min", BLUE),
    ("naming → 50 x402 micropayments → screened $5 escrow", None),
    ("delivery → watcher auto-release → dispute / refund branch", None),
    ("every path balance-asserted on-chain", None),
]
nxt = [
    "Deploy to Arc Testnet — artifact-driven, one command",
    "Activate Circle agent wallets with the API key",
    "Circle Compliance Engine live screening",
    "48h unattended soak on real rails",
]
tb, tf = textbox(s, MX, Inches(2.75), CW / 2 - Inches(0.2), Inches(3.5))
p = tf.paragraphs[0]; _set_run(p.add_run(), "WORKING NOW", 12, LOW, MONO, spacing=1.4)
for text, code in now:
    p = tf.add_paragraph(); p.space_before = Pt(11); p.line_spacing = 1.3
    _set_run(p.add_run(), "✓  ", 14, JADE, MONO, bold=True)
    _set_run(p.add_run(), text, 14, MID if code is None else HI, SANS)

tb, tf = textbox(s, MX + CW / 2, Inches(2.75), CW / 2 - Inches(0.2), Inches(3.5))
p = tf.paragraphs[0]; _set_run(p.add_run(), "NEXT", 12, LOW, MONO, spacing=1.4)
for text in nxt:
    p = tf.add_paragraph(); p.space_before = Pt(11); p.line_spacing = 1.3
    _set_run(p.add_run(), "→  ", 14, GOLD, MONO, bold=True)
    _set_run(p.add_run(), text, 14, MID, SANS)


# =====================================================================
# 11 · CLOSE
# =====================================================================
s = slide()
eyebrow(s, "10", "The thesis")
tb, tf = textbox(s, MX, Inches(2.0), CW, Inches(2.2))
p = tf.paragraphs[0]; p.line_spacing = 1.12
_set_run(p.add_run(), "When agents transact at machine speed,", 32, HI, MONO, bold=True, spacing=-0.4)
p = tf.add_paragraph(); p.line_spacing = 1.12
_set_run(p.add_run(), "settlement has to be ", 32, HI, MONO, bold=True, spacing=-0.4)
_set_run(p.add_run(), "named, instant, and", 32, JADE, MONO, bold=True, spacing=-0.4)
p = tf.add_paragraph(); p.line_spacing = 1.12
_set_run(p.add_run(), "compliant", 32, JADE, MONO, bold=True, spacing=-0.4)
_set_run(p.add_run(), " — or it doesn't happen at all.", 32, HI, MONO, bold=True, spacing=-0.4)

lead(s, Inches(4.5), [
    ("AgentMesh", HI, True),
    (" is that settlement layer: identity, sub-cent payments, screened escrow, and automation — stablecoin-native on Arc, ready for the agent economy.", MID, False),
], width=Inches(11), size=17)

pill = rrect(s, MX, Inches(6.25), Inches(9.2), Inches(0.6), fill=PANEL2, line=JADE)
pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
pp = pill.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
_set_run(pp.add_run(), "●  AgentMesh · Programmable Money Hackathon · Agentic Economy", 13, JADE, MONO, spacing=0.5)


import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "AgentMesh-deck.pptx")
prs.save(out)
print("saved", out)
